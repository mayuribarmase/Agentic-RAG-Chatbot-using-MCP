# ingestion_agent.py

import fitz  # PyMuPDF
import docx
import pptx
import pandas as pd
import os

def parse_pdf(filepath):
    text = ""
    with fitz.open(filepath) as doc:
        for page in doc:
            text += page.get_text()
    return text

def parse_docx(filepath):
    doc = docx.Document(filepath)
    return "\n".join([p.text for p in doc.paragraphs])

def parse_pptx(filepath):
    prs = pptx.Presentation(filepath)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def parse_csv(filepath):
    df = pd.read_csv(filepath)
    return df.to_string()

def parse_txt(filepath):
    with open(filepath, "r", encoding="utf-8") as f:
        return f.read()

def parse_md(filepath):
    with open(filepath, "r", encoding="utf-8") as f:
        return f.read()

def ingest_document(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".pdf":
        return parse_pdf(filepath)
    elif ext == ".docx":
        return parse_docx(filepath)
    elif ext in [".pptx"]:
        return parse_pptx(filepath)
    elif ext == ".csv":
        return parse_csv(filepath)
    elif ext == ".md":
        return parse_md(filepath)
    else:  # txt fallback
        return parse_txt(filepath)
